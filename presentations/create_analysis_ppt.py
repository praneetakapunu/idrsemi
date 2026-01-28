#!/usr/bin/env python3
"""
Create IDRSemi Strategic Analysis Presentation
Comparing Pacetronix Partnership vs RISC-V Development
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from io import BytesIO
import numpy as np

# Color scheme - Professional blue/teal theme
PRIMARY_COLOR = RGBColor(0, 48, 135)      # Deep blue
ACCENT_COLOR = RGBColor(0, 176, 240)      # Bright blue
SUCCESS_COLOR = RGBColor(34, 177, 76)     # Green
WARNING_COLOR = RGBColor(255, 127, 39)    # Orange
DANGER_COLOR = RGBColor(237, 28, 36)      # Red
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(191, 191, 191)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

def add_agenda_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Agenda"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    agenda_items = [
        "Executive Summary",
        "Market Analysis",
        "Option 1: Pacetronix Partnership",
        "Option 2: RISC-V Development", 
        "Financial Projections",
        "Risk Assessment",
        "Recommendation & Next Steps"
    ]
    
    for i, item in enumerate(agenda_items):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_after = Pt(12)

def add_executive_summary(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add key decision box
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(1.2)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = ACCENT_COLOR
    
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Strategic Decision Point"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Choose between immediate revenue (Pacetronix) or long-term IP ownership (RISC-V)"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add comparison points
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    tf.margin_top = Inches(2.5)
    
    comparisons = [
        ("Timeline", "Pacetronix: 6-9 months to revenue | RISC-V: 12-18 months"),
        ("Investment", "Pacetronix: $50K-100K | RISC-V: $200K-500K"),
        ("Risk Level", "Pacetronix: Medium (regulatory) | RISC-V: High (market)"),
        ("Growth Potential", "Pacetronix: Steady | RISC-V: Exponential")
    ]
    
    for label, comparison in comparisons:
        p = tf.add_paragraph()
        p.text = f"• {label}: {comparison}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

def add_market_analysis(prs):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Analysis"
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    
    # Create market size comparison chart
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    fig.patch.set_facecolor('white')
    
    # Medical device market
    categories1 = ['Pacemakers', 'Defibrillators', 'Monitoring', 'Other']
    sizes1 = [35, 25, 25, 15]
    colors1 = ['#0030B7', '#00B0F0', '#22B14C', '#FF7F27']
    
    ax1.pie(sizes1, labels=categories1, colors=colors1, autopct='%1.1f%%', startangle=90)
    ax1.set_title('Medical Device IC Market\n$4.2B Total', fontsize=14, weight='bold')
    
    # RISC-V market
    categories2 = ['IoT', 'Automotive', 'Data Center', 'Consumer', 'Industrial']
    sizes2 = [30, 25, 20, 15, 10]
    colors2 = ['#0030B7', '#00B0F0', '#22B14C', '#FF7F27', '#ED1C24']
    
    ax2.pie(sizes2, labels=categories2, colors=colors2, autopct='%1.1f%%', startangle=90)
    ax2.set_title('RISC-V Market Segments\n$2.7B by 2027', fontsize=14, weight='bold')
    
    plt.tight_layout()
    
    # Save and add to slide
    img_stream = BytesIO()
    plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))
    
    # Add growth projections
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Market Insights:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY_COLOR
    
    insights = [
        "Medical device market: 7.2% CAGR, highly regulated but stable",
        "RISC-V market: 35% CAGR, emerging technology with high growth potential",
        "Both markets offer significant opportunities with different risk profiles"
    ]
    
    for insight in insights:
        p = tf.add_paragraph()
        p.text = f"• {insight}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)

def add_pacetronix_option(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Option 1: Pacetronix Partnership"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Two column layout
    left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    right_col = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    
    # Pros
    tf_left = left_col.text_frame
    p = tf_left.paragraphs[0]
    p.text = "Advantages"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = SUCCESS_COLOR
    
    pros = [
        "Immediate revenue opportunity",
        "Established customer with clear needs",
        "High-margin medical market",
        "Valuable regulatory experience",
        "Potential for long-term partnership",
        "Reference customer for future deals"
    ]
    
    for pro in pros:
        p = tf_left.add_paragraph()
        p.text = f"✓ {pro}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Cons
    tf_right = right_col.text_frame
    p = tf_right.paragraphs[0]
    p.text = "Challenges"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = WARNING_COLOR
    
    cons = [
        "Strict FDA compliance required",
        "18-24 month development cycle",
        "High liability insurance needs",
        "Conservative design constraints",
        "Limited IP ownership",
        "Dependency on single customer"
    ]
    
    for con in cons:
        p = tf_right.add_paragraph()
        p.text = f"⚠ {con}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

def add_technical_requirements_pacetronix(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Pacetronix: Technical Requirements"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Main requirements
    p = tf.paragraphs[0]
    p.text = "Core Specifications"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    
    specs = [
        ("Process Node", "180nm/130nm (mature, reliable)"),
        ("Power Budget", "< 10µW active, < 100nA standby"),
        ("Operating Voltage", "1.0V - 3.6V (battery operation)"),
        ("Temperature Range", "-10°C to +55°C (implant conditions)"),
        ("Lifetime", "10+ years continuous operation"),
        ("Safety", "Triple redundancy, fail-safe mechanisms")
    ]
    
    for spec, detail in specs:
        p = tf.add_paragraph()
        p.text = f"• {spec}: {detail}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # DV Focus Box
    left = Inches(1)
    top = Inches(5.2)
    width = Inches(8)
    height = Inches(1.8)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = ACCENT_COLOR
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "DV Excellence Required"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "• ISO 13485 compliant verification\n• 100% code coverage mandatory\n• Formal verification for critical blocks"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.LEFT

def add_riscv_option(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Option 2: RISC-V Development"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Two column layout
    left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    right_col = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    
    # Pros
    tf_left = left_col.text_frame
    p = tf_left.paragraphs[0]
    p.text = "Advantages"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = SUCCESS_COLOR
    
    pros = [
        "Full IP ownership",
        "Hot technology trend",
        "Large ecosystem support",
        "Multiple market opportunities",
        "Design freedom & innovation",
        "Potential for licensing revenue"
    ]
    
    for pro in pros:
        p = tf_left.add_paragraph()
        p.text = f"✓ {pro}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Cons
    tf_right = right_col.text_frame
    p = tf_right.paragraphs[0]
    p.text = "Challenges"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = WARNING_COLOR
    
    cons = [
        "No immediate customer",
        "Highly competitive market",
        "Longer development timeline",
        "Higher upfront investment",
        "Marketing & sales required",
        "Need strong differentiation"
    ]
    
    for con in cons:
        p = tf_right.add_paragraph()
        p.text = f"⚠ {con}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

def add_riscv_technical_approach(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "RISC-V: Technical Approach"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Architecture choices
    p = tf.paragraphs[0]
    p.text = "Proposed Architecture: RV32IMC + Custom Extensions"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT_COLOR
    
    features = [
        "Base: RV32I (minimal 32-bit integer ISA)",
        "M Extension: Hardware multiply/divide",
        "C Extension: Compressed instructions (code density)",
        "Custom: Domain-specific accelerators"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    # Differentiation box
    p = tf.add_paragraph()
    p.text = "\nKey Differentiators"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = SUCCESS_COLOR
    p.space_before = Pt(12)
    
    differentiators = [
        ("Security Focus", "Hardware root of trust, secure boot, crypto extensions"),
        ("Power Optimization", "Dynamic voltage scaling, power domains, sleep modes"),
        ("AI/ML Ready", "Custom instructions for inference, SIMD operations"),
        ("Verification IP", "Pre-verified UVM testbenches included")
    ]
    
    for diff, desc in differentiators:
        p = tf.add_paragraph()
        p.text = f"• {diff}: {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

def add_financial_projections(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Projections (3-Year)"
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    
    # Create financial comparison chart
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor('white')
    
    years = ['Year 1', 'Year 2', 'Year 3']
    pacetronix_revenue = [0.5, 1.2, 2.0]  # Millions
    riscv_revenue = [0, 0.8, 3.5]  # Millions
    
    x = np.arange(len(years))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, pacetronix_revenue, width, label='Pacetronix Path', color='#00B0F0')
    bars2 = ax.bar(x + width/2, riscv_revenue, width, label='RISC-V Path', color='#22B14C')
    
    ax.set_ylabel('Revenue ($M)', fontsize=12, weight='bold')
    ax.set_xlabel('Timeline', fontsize=12, weight='bold')
    ax.set_title('Revenue Projection Comparison', fontsize=16, weight='bold', pad=20)
    ax.set_xticks(x)
    ax.set_xticklabels(years)
    ax.legend(loc='upper left', fontsize=12)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    # Add value labels on bars
    for bars in [bars1, bars2]:
        for bar in bars:
            height = bar.get_height()
            if height > 0:
                ax.annotate(f'${height}M',
                           xy=(bar.get_x() + bar.get_width() / 2, height),
                           xytext=(0, 3),
                           textcoords="offset points",
                           ha='center', va='bottom',
                           fontsize=10, weight='bold')
    
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3), width=Inches(9))
    
    # Add financial details table
    table_data = [
        ["Metric", "Pacetronix Path", "RISC-V Path"],
        ["Initial Investment", "$50K - $100K", "$200K - $500K"],
        ["Break-even", "Month 9", "Month 18"],
        ["3-Year Revenue", "$3.7M", "$4.3M"],
        ["Gross Margin", "45%", "65%"],
        ["Risk Level", "Medium", "High"]
    ]
    
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(1.5)
    top = Inches(5.2)
    width = Inches(7)
    height = Inches(1.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Format table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            
            if i == 0:  # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
            else:
                if j == 0:  # First column
                    paragraph.font.bold = True
                paragraph.font.color.rgb = DARK_GRAY

def add_risk_assessment(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Risk Assessment Matrix"
    p.font.size = Pt(36)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    
    # Create risk matrix visualization
    fig, ax = plt.subplots(figsize=(10, 6))
    fig.patch.set_facecolor('white')
    
    # Risk data (Impact, Probability) for each risk
    pacetronix_risks = {
        'Regulatory Delays': (4, 3),
        'Customer Dependency': (3, 2),
        'Technical Complexity': (2, 2),
        'Liability Issues': (5, 1),
    }
    
    riscv_risks = {
        'Market Competition': (3, 4),
        'No Initial Customer': (4, 5),
        'Technology Adoption': (3, 3),
        'Funding Gap': (4, 2),
    }
    
    # Plot risk matrix background
    for i in range(5):
        for j in range(5):
            if i * j < 6:
                color = '#90EE90'  # Light green (low risk)
            elif i * j < 12:
                color = '#FFD700'  # Gold (medium risk)
            else:
                color = '#FF6B6B'  # Light red (high risk)
            ax.add_patch(plt.Rectangle((i, j), 1, 1, facecolor=color, alpha=0.3, edgecolor='gray'))
    
    # Plot risks
    for risk, (impact, prob) in pacetronix_risks.items():
        ax.scatter(prob-0.5, impact-0.5, s=300, c='#00B0F0', marker='o', edgecolors='black', linewidths=2)
        ax.annotate(risk, (prob-0.5, impact-0.5), xytext=(5, 5), textcoords='offset points', 
                   fontsize=9, ha='left', bbox=dict(boxstyle='round,pad=0.3', facecolor='white', alpha=0.7))
    
    for risk, (impact, prob) in riscv_risks.items():
        ax.scatter(prob-0.5, impact-0.5, s=300, c='#22B14C', marker='s', edgecolors='black', linewidths=2)
        ax.annotate(risk, (prob-0.5, impact-0.5), xytext=(5, -15), textcoords='offset points',
                   fontsize=9, ha='left', bbox=dict(boxstyle='round,pad=0.3', facecolor='white', alpha=0.7))
    
    ax.set_xlim(0, 5)
    ax.set_ylim(0, 5)
    ax.set_xlabel('Probability →', fontsize=14, weight='bold')
    ax.set_ylabel('Impact →', fontsize=14, weight='bold')
    ax.set_title('Risk Matrix: Pacetronix (●) vs RISC-V (■)', fontsize=16, weight='bold', pad=20)
    ax.grid(True, linestyle='-', alpha=0.5)
    
    # Set tick labels
    ax.set_xticks([0.5, 1.5, 2.5, 3.5, 4.5])
    ax.set_xticklabels(['Very Low', 'Low', 'Medium', 'High', 'Very High'])
    ax.set_yticks([0.5, 1.5, 2.5, 3.5, 4.5])
    ax.set_yticklabels(['Very Low', 'Low', 'Medium', 'High', 'Very High'])
    
    plt.tight_layout()
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='PNG', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9))

def add_recommendation_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Strategic Recommendation"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Main recommendation box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1.5)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SUCCESS_COLOR
    shape.line.color.rgb = SUCCESS_COLOR
    
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Recommended Strategy: Hybrid Approach"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Start with Pacetronix for immediate revenue while developing RISC-V IP in parallel"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Implementation timeline
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    tf.margin_top = Inches(2.2)
    
    p = tf.paragraphs[0]
    p.text = "Implementation Timeline"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    
    timeline = [
        ("Months 1-3", "Finalize Pacetronix contract, Begin architecture design"),
        ("Months 4-6", "Medical ASIC development, Start RISC-V core design"),
        ("Months 7-9", "Pacetronix verification phase, RISC-V prototype"),
        ("Months 10-12", "Pacetronix tape-out, RISC-V customer demos"),
        ("Year 2", "Pacetronix production, RISC-V first customers")
    ]
    
    for period, activities in timeline:
        p = tf.add_paragraph()
        p.text = f"• {period}: {activities}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

def add_next_steps_slide(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Immediate Next Steps"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Pacetronix actions
    p = tf.paragraphs[0]
    p.text = "Pacetronix Partnership (This Week)"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    
    pacetronix_steps = [
        "Schedule technical deep-dive meeting",
        "Request detailed specifications",
        "Discuss NDA and contract terms",
        "Estimate resource requirements",
        "Define success metrics"
    ]
    
    for step in pacetronix_steps:
        p = tf.add_paragraph()
        p.text = f"✓ {step}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    # RISC-V preparation
    p = tf.add_paragraph()
    p.text = "\nRISC-V Preparation (Next 30 Days)"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    p.space_before = Pt(12)
    
    riscv_steps = [
        "Complete market research & competitive analysis",
        "Define core architecture & differentiators",
        "Build DV infrastructure & methodology",
        "Identify potential pilot customers",
        "Create funding plan & investor pitch"
    ]
    
    for step in riscv_steps:
        p = tf.add_paragraph()
        p.text = f"◆ {step}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

def add_appendix_dv_focus(prs):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Appendix: DV Strategy for Both Paths"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Leveraging Your DV Expertise"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    
    # DV approach
    dv_items = [
        ("Pacetronix DV Focus", [
            "- ISO 13485 compliant verification methodology",
            "- Fault injection & safety verification",
            "- Power-aware verification for battery life",
            "- Analog mixed-signal verification"
        ]),
        ("RISC-V DV Advantage", [
            "- Pre-built UVM testbenches as product differentiator",
            "- RISC-V compliance test suite integration",
            "- Performance verification infrastructure",
            "- Open-source contribution opportunity"
        ]),
        ("Common DV Infrastructure", [
            "- Reusable verification IP development",
            "- Continuous integration setup",
            "- Coverage-driven verification",
            "- Formal verification for critical blocks"
        ])
    ]
    
    for category, items in dv_items:
        p = tf.add_paragraph()
        p.text = f"\n{category}:"
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = PRIMARY_COLOR
        p.space_before = Pt(8)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)

# Main presentation creation
def create_presentation():
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add all slides
    add_title_slide(prs, 
                   "IDRSemi Strategic Analysis",
                   "Pacetronix Partnership vs RISC-V Development\nJanuary 2026")
    
    add_agenda_slide(prs)
    add_executive_summary(prs)
    add_market_analysis(prs)
    add_pacetronix_option(prs)
    add_technical_requirements_pacetronix(prs)
    add_riscv_option(prs)
    add_riscv_technical_approach(prs)
    add_financial_projections(prs)
    add_risk_assessment(prs)
    add_recommendation_slide(prs)
    add_next_steps_slide(prs)
    add_appendix_dv_focus(prs)
    
    # Save presentation
    prs.save('/home/exedev/clawd/idrsemi/presentations/IDRSemi_Strategic_Analysis.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()